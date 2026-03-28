from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "Provincial Health Control Center", "Executive Dashboard Redesign Overview")
    
    # Slide 2: The New Architecture
    add_bullet_slide(prs, "The New Architecture", [
        "A fundamentally transformed 3-Tab Control Center.",
        "Replaces static reports with dynamic, interactive data environments.",
        "Built-in 'Synthesis Engine' mathematically identifies priority intervention targets.",
        "Maintains a clinical, high-contrast aesthetic design."
    ])
    
    # Slide 3: Tab 1 (Equity)
    add_bullet_slide(prs, "Tab 1: Population Health Equity (🌎)", [
        "Tracks Community Health Service Areas (CHSAs) for access gaps.",
        "Alert Banner: Automatically flags the 'Highest Priority Community' using a custom vulnerability score.",
        "Access Gaps: Visualizes % without a Family Doctor by region.",
        "Wealth-Health Gap: Scatter plot correlating income against life expectancy.",
        "Opioid Vulnerability: Clean blue/green heatmap highlighting top 10 impacted regions."
    ])
    
    # Slide 4: Tab 2 (Wait Times)
    add_bullet_slide(prs, "Tab 2: BC Surgical Wait Times (⏱️)", [
        "Dedicated tracker for measuring surgical efficacy across 8 procedures.",
        "Historical Comparison: Generates direct procedure volume and wait comparisons against a 2014 baseline.",
        "Trend Matrix: Visualizes BC performance mapped over top of the National Average.",
        "Current Performance Bar: Visually categorizes all procedures into Red/Amber/Green based on Federal Benchmark targets."
    ])
    
    # Slide 5: Tab 3 (Opioids)
    add_bullet_slide(prs, "Tab 3: Opioid Crisis Deep Dive (⚠️)", [
        "Acute crisis management environment.",
        "Narrative Callout: Explicitly links surging toxicity metrics directly to the most vulnerable local community identified in Tab 1.",
        "YoY Deltas: Tracks year-over-year shifts in toxicity deaths, hospitalizations, and ER visits.",
        "Dual-Axis Plot: Matches absolute toxicity deaths against the sheer volume of structural hospital load.",
        "Inter-Provincial Ranking: Standardized death rate comparisons against the rest of Canada."
    ])
    
    # Slide 6: Deployment
    add_bullet_slide(prs, "Next Steps & Deployment", [
        "The Control Center is fully functional and successfully tested locally.",
        "Interactive navigation automatically prevents sidebar filter conflicts between tabs.",
        "Codebase has been formatted for instant push to Streamlit Cloud.",
        "Ready to be deployed out to Population Health Analysts across the Authority."
    ])
    
    # Save the file
    prs.save("BC_Health_Dashboard_Presentation.pptx")
    print("Presentation created successfully: BC_Health_Dashboard_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
